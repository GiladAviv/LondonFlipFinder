"""Render a python-pptx-authored deck to PDF without LibreOffice.

The decks in this repo are built entirely from rectangles, text boxes, tables and
pictures at absolute positions, so every shape maps cleanly onto a reportlab
canvas.  Fonts are metric-compatible substitutes: Gelasio for Georgia, Carlito
for Calibri, DejaVu Sans Mono for Consolas.
"""

import io
import os
import sys

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from reportlab.lib.utils import ImageReader
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfgen import canvas as rl_canvas

EMU_PT = 12700.0
HERE = os.path.dirname(os.path.abspath(__file__))
FONT_DIR = os.path.join(HERE, "fonts")

_FACES = {
    "Georgia": ("Gelasio", "Gelasio-Regular.ttf", "Gelasio-Bold.ttf",
                "Gelasio-RegularItalic.ttf", "Gelasio-BoldItalic.ttf"),
    "Calibri": ("Carlito", "Carlito-Regular.ttf", "Carlito-Bold.ttf",
                "Carlito-Italic.ttf", "Carlito-Italic.ttf"),
}
_MONO = "/usr/share/fonts/dejavu-sans-mono-fonts/DejaVuSansMono.ttf"
_MONO_B = "/usr/share/fonts/dejavu-sans-mono-fonts/DejaVuSansMono-Bold.ttf"


def register_fonts():
    for family, (base, reg, bold, ital, bital) in _FACES.items():
        pdfmetrics.registerFont(TTFont(base, os.path.join(FONT_DIR, reg)))
        pdfmetrics.registerFont(TTFont(base + "-Bold", os.path.join(FONT_DIR, bold)))
        pdfmetrics.registerFont(TTFont(base + "-It", os.path.join(FONT_DIR, ital)))
        pdfmetrics.registerFont(TTFont(base + "-BoldIt", os.path.join(FONT_DIR, bital)))
        pdfmetrics.registerFontFamily(base, normal=base, bold=base + "-Bold",
                                      italic=base + "-It", boldItalic=base + "-BoldIt")
    pdfmetrics.registerFont(TTFont("Mono", _MONO))
    pdfmetrics.registerFont(TTFont("Mono-Bold", _MONO_B))


def font_for(name, bold, italic):
    base = {"Georgia": "Gelasio", "Calibri": "Carlito"}.get(name)
    if base is None:
        return "Mono-Bold" if bold else "Mono"
    if bold and italic:
        return base + "-BoldIt"
    if bold:
        return base + "-Bold"
    if italic:
        return base + "-It"
    return base


def rgb_of(color, default=(0, 0, 0)):
    try:
        if color is None or color.type is None:
            return default
        v = color.rgb
        return (v[0] / 255.0, v[1] / 255.0, v[2] / 255.0)
    except Exception:
        return default


# --------------------------------------------------------------------------- text


class Run:
    __slots__ = ("text", "font", "size", "color")

    def __init__(self, text, font, size, color):
        self.text = text
        self.font = font
        self.size = size
        self.color = color

    def width(self):
        return pdfmetrics.stringWidth(self.text, self.font, self.size)


def runs_of(paragraph, defaults):
    d_size, d_color, d_font, d_bold = defaults
    out = []
    for r in paragraph.runs:
        f = r.font
        size = f.size.pt if f.size is not None else d_size
        bold = d_bold if f.bold is None else f.bold
        italic = bool(f.italic)
        name = f.name or d_font
        out.append(Run(r.text, font_for(name, bold, italic), size,
                       rgb_of(f.color, d_color)))
    return out


def wrap_runs(runs, max_width):
    """Greedy word wrap over a run sequence; returns list of lines of runs."""
    lines = [[]]
    width = 0.0
    for run in runs:
        # split keeping spaces attached to the preceding token
        tokens, buf = [], ""
        for ch in run.text:
            buf += ch
            if ch == " ":
                tokens.append(buf)
                buf = ""
        if buf:
            tokens.append(buf)
        for tok in tokens:
            w = pdfmetrics.stringWidth(tok, run.font, run.size)
            if width + w > max_width and lines[-1] and tok.strip():
                stripped = pdfmetrics.stringWidth(tok.rstrip(), run.font, run.size)
                if width + stripped > max_width:
                    lines.append([])
                    width = 0.0
            if lines[-1] and lines[-1][-1].font == run.font and \
                    lines[-1][-1].size == run.size and lines[-1][-1].color == run.color:
                lines[-1][-1].text += tok
            else:
                lines[-1].append(Run(tok, run.font, run.size, run.color))
            width += w
    return lines


def layout_text(tf, box_w, defaults):
    """Return (blocks, total_height) where each block is (lines, lh, align, space_after)."""
    blocks = []
    total = 0.0
    for para in tf.paragraphs:
        runs = runs_of(para, defaults)
        sizes = [r.size for r in runs] or [defaults[0]]
        size = max(sizes)
        spc = para.line_spacing if isinstance(para.line_spacing, float) else \
            (para.line_spacing if para.line_spacing else 1.0)
        if hasattr(spc, "pt"):
            lh = spc.pt
        else:
            lh = size * 1.2 * float(spc)
        after = para.space_after.pt if para.space_after is not None else 0.0
        lines = wrap_runs(runs, box_w) if runs else [[]]
        align = para.alignment
        blocks.append((lines, lh, align, after, size))
        total += lh * len(lines) + after
    return blocks, total


def draw_text_frame(c, tf, x, y_top, w, h, defaults, anchor=None):
    ml, mr = tf.margin_left.pt, tf.margin_right.pt
    mt, mb = tf.margin_top.pt, tf.margin_bottom.pt
    bx, bw = x + ml, w - ml - mr
    if bw <= 1:
        return
    blocks, total = layout_text(tf, bw, defaults)
    anchor = anchor if anchor is not None else tf.vertical_anchor
    top = y_top - mt
    if anchor == MSO_ANCHOR.MIDDLE:
        avail = h - mt - mb
        top = y_top - mt - max(0.0, (avail - total) / 2.0)
    elif anchor == MSO_ANCHOR.BOTTOM:
        top = y_top - h + mb + total

    cur = top
    for lines, lh, align, after, size in blocks:
        for line in lines:
            line_w = sum(r.width() for r in line)
            if align == PP_ALIGN.CENTER:
                lx = bx + (bw - line_w) / 2.0
            elif align == PP_ALIGN.RIGHT:
                lx = bx + bw - line_w
            else:
                lx = bx
            ascent = pdfmetrics.getAscent(line[0].font, size) if line else size * 0.75
            baseline = cur - (lh - size) / 2.0 - ascent
            for r in line:
                if r.text.strip():
                    c.setFillColorRGB(*r.color)
                    c.setFont(r.font, r.size)
                    c.drawString(lx, baseline, r.text)
                lx += r.width()
            cur -= lh
        cur -= after


# -------------------------------------------------------------------------- shapes


def draw_table(c, gf, H):
    tbl = gf.table
    col_w = [col.width / EMU_PT for col in tbl.columns]
    row_h = [row.height / EMU_PT for row in tbl.rows]
    x0, y0 = gf.left / EMU_PT, gf.top / EMU_PT
    # PowerPoint grows rows to fit content; scale row heights to the frame height
    frame_h = gf.height / EMU_PT
    if sum(row_h) > 0 and abs(sum(row_h) - frame_h) > 1:
        k = frame_h / sum(row_h)
        row_h = [r * k for r in row_h]
    y = y0
    for ri, row in enumerate(tbl.rows):
        x = x0
        for ci, cell in enumerate(row.cells):
            w, h = col_w[ci], row_h[ri]
            fill = cell.fill
            try:
                if str(fill.type).startswith("SOLID"):
                    c.setFillColorRGB(*rgb_of(fill.fore_color, (1, 1, 1)))
                    c.rect(x, H - y - h, w, h, fill=1, stroke=0)
            except Exception:
                pass
            draw_text_frame(c, cell.text_frame, x + cell.margin_left.pt,
                            H - y - cell.margin_top.pt,
                            w - cell.margin_left.pt - cell.margin_right.pt,
                            h - cell.margin_top.pt - cell.margin_bottom.pt,
                            (11.0, (0.1, 0.1, 0.1), "Calibri", False),
                            anchor=cell.vertical_anchor)
            x += w
        y += row_h[ri]


_A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
_R_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"


def _poster_blob(shape):
    """The poster-frame image behind a movie shape, via its blipFill relationship."""
    # The blipFill on a movie sits in the presentationml namespace, so search for the blip
    # itself rather than assuming its parent's namespace.
    blip = shape._element.find(f".//{_A_NS}blip")
    if blip is None:
        return None
    rid = blip.get(f"{_R_NS}embed")
    if not rid:
        return None
    return shape.part.related_part(rid).blob


def draw_shape(c, sh, H):
    st = sh.shape_type
    if st == MSO_SHAPE_TYPE.GROUP:
        for sub in sh.shapes:
            draw_shape(c, sub, H)
        return
    if getattr(sh, "has_table", False):
        draw_table(c, sh, H)
        return
    x = sh.left / EMU_PT
    y = sh.top / EMU_PT
    w = sh.width / EMU_PT
    h = sh.height / EMU_PT
    if st == MSO_SHAPE_TYPE.PICTURE:
        img = ImageReader(io.BytesIO(sh.image.blob))
        c.drawImage(img, x, H - y - h, w, h, mask="auto")
        return
    if st == MSO_SHAPE_TYPE.MEDIA:
        # A movie has no .image; its poster frame is the blipFill on the same p:pic element.
        blob = _poster_blob(sh)
        if blob is not None:
            c.drawImage(ImageReader(io.BytesIO(blob)), x, H - y - h, w, h, mask="auto")
        return
    if st == MSO_SHAPE_TYPE.AUTO_SHAPE:
        fill = sh.fill
        try:
            if str(fill.type).startswith("SOLID"):
                c.setFillColorRGB(*rgb_of(fill.fore_color, (1, 1, 1)))
                c.rect(x, H - y - h, w, h, fill=1, stroke=0)
            elif str(fill.type).startswith("GRADIENT"):
                # Approximate a linear gradient with horizontal bands so the preview
                # reflects what PowerPoint will actually show.
                stops = list(fill.gradient_stops)
                c0 = rgb_of(stops[0].color, (1, 1, 1))
                c1 = rgb_of(stops[-1].color, (1, 1, 1))
                bands = 120
                for i in range(bands):
                    t = i / (bands - 1.0)
                    c.setFillColorRGB(*[c0[k] + (c1[k] - c0[k]) * t for k in range(3)])
                    bh = h / bands + 0.6
                    c.rect(x, H - y - h + (h * (1 - t)) - bh, w, bh, fill=1, stroke=0)
        except Exception:
            pass
    if sh.has_text_frame and sh.text_frame.text.strip():
        draw_text_frame(c, sh.text_frame, x, H - y, w, h,
                        (18.0, (0.05, 0.12, 0.2), "Calibri", False))


def convert(src, dst):
    prs = Presentation(src)
    W = prs.slide_width / EMU_PT
    H = prs.slide_height / EMU_PT
    c = rl_canvas.Canvas(dst, pagesize=(W, H))
    c.setTitle(os.path.splitext(os.path.basename(dst))[0].replace("_", " "))
    for i, slide in enumerate(prs.slides, 1):
        c.setFillColorRGB(1, 1, 1)
        c.rect(0, 0, W, H, fill=1, stroke=0)
        for sh in slide.shapes:
            draw_shape(c, sh, H)
        c.bookmarkPage("slide%d" % i)
        c.addOutlineEntry("Slide %d" % i, "slide%d" % i, level=0)
        c.showPage()
    c.save()
    return len(prs.slides._sldIdLst)


if __name__ == "__main__":
    register_fonts()
    convert(sys.argv[1], sys.argv[2])
    print("wrote", sys.argv[2], os.path.getsize(sys.argv[2]), "bytes")
