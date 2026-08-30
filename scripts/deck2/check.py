# -*- coding: utf-8 -*-
"""Check presentation.pptx against the brief's hard clarity limits."""
import re, sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt

ROOT = Path(__file__).resolve().parents[2]
sys.path.insert(0, str(Path(__file__).resolve().parent))
from deck_content import MAIN, APPENDIX

MAX_BULLETS, MAX_WORDS_BULLET, MAX_WORDS_SLIDE, MIN_PT = 5, 12, 42, 16
bad = []

def words(t):
    return len(re.sub(r"\*\*", "", t).split())

for i, spec in enumerate(MAIN + APPENDIX, 1):
    tag = f"slide {i}"
    bl = spec.get("bullets", [])
    if len(bl) > MAX_BULLETS:
        bad.append(f"{tag}: {len(bl)} bullets (max {MAX_BULLETS})")
    for b in bl:
        if words(b) > MAX_WORDS_BULLET:
            bad.append(f"{tag}: bullet {words(b)} words > {MAX_WORDS_BULLET}: {b!r}")
    tot = sum(words(b) for b in bl)
    if tot > MAX_WORDS_SLIDE:
        bad.append(f"{tag}: {tot} bullet words > {MAX_WORDS_SLIDE}")
    emph = sum(len(re.findall(r"\*\*(.+?)\*\*", b)) for b in bl) + \
           sum(len(re.findall(r"\*\*(.+?)\*\*", c)) for c in spec.get("code", []))
    if emph > 1:
        bad.append(f"{tag}: {emph} emphasised spans (max 1)")
    if len(spec.get("code", [])) > 8:
        bad.append(f"{tag}: {len(spec['code'])} code lines > 8")
    t = spec.get("table")
    if t:
        if len(t["rows"]) + 1 > 6 and not spec.get("dense"):
            bad.append(f"{tag}: table {len(t['rows'])+1} rows > 6")
        if len(t["headers"]) > 5:
            bad.append(f"{tag}: table {len(t['headers'])} cols > 5")
        if abs(sum(t["widths"]) - sum(t["widths"])) > 0:
            pass
        if sum(t["widths"]) > 11.9:
            bad.append(f"{tag}: table width {sum(t['widths']):.2f}in > 11.9")
    ns = len(re.findall(r"[.!?] ", spec["notes"])) + 1
    if ns < 5:
        bad.append(f"{tag}: only ~{ns} sentences of notes (min 5)")
    if "Source: notebook" not in spec["notes"] and spec["layout"] != "diagram":
        if "cell" not in spec["notes"]:
            bad.append(f"{tag}: notes give no cell numbers")

prs = Presentation(str(ROOT / "presentation.pptx"))
for n, sl in enumerate(prs.slides, 1):
    if not sl.has_notes_slide or not sl.notes_slide.notes_text_frame.text.strip():
        bad.append(f"slide {n}: no speaker notes")
    for sh in sl.shapes:
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.size and r.font.size < Pt(MIN_PT) and r.text.strip():
                        bad.append(f"slide {n}: {r.font.size.pt}pt text {r.text[:35]!r}")
        if getattr(sh, "has_table", False):
            for row in sh.table.rows:
                for c in row.cells:
                    for p in c.text_frame.paragraphs:
                        for r in p.runs:
                            if r.font.size and r.font.size < Pt(MIN_PT) and r.text.strip():
                                bad.append(f"slide {n}: TABLE {r.font.size.pt}pt {r.text[:25]!r}")

print(f"slides: {len(prs.slides)}  ({len(MAIN)} main + {len(APPENDIX)} appendix)")
if bad:
    print(f"\n{len(bad)} VIOLATIONS:")
    for b in bad:
        print("  -", b)
else:
    print("\nAll clarity limits pass.")
