# -*- coding: utf-8 -*-
"""Build presentation.pptx from deck_content.py with python-pptx."""
from __future__ import annotations

import re
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
sys.path.insert(0, str(Path(__file__).resolve().parent))
from deck_content import APPENDIX, MAIN  # noqa: E402

# ------------------------------------------------------------------ design system
ACCENT = RGBColor(0x2A, 0x78, 0xD6)     # the notebook's own primary series blue
DEEP = RGBColor(0x10, 0x42, 0x81)       # deep blue, titles
INK = RGBColor(0x2B, 0x2B, 0x2B)        # body text
MUTED = RGBColor(0x6B, 0x6A, 0x67)
RULE = RGBColor(0xD9, 0xD8, 0xD4)
PANEL = RGBColor(0xF3, 0xF6, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PAGE = RGBColor(0x9A, 0x99, 0x95)
BG_TOP = RGBColor(0xFF, 0xFF, 0xFF)     # gradient wash: near-white top-left
BG_BOT = RGBColor(0xE8, 0xF0, 0xF9)     # to pale blue bottom-right

BODY_FONT = "Calibri"
MONO_FONT = "Consolas"

SW, SH = 13.333, 7.5
ML, MR = 0.75, 0.75
CW = SW - ML - MR                        # 11.833 content width

T_KICKER, T_TITLE, T_BULLET, T_SUB = 16, 30, 19, 17
T_TABLE, T_TABLE_HEAD, T_CODE, T_PAGE = 18, 18, 17, 16

CONTENT_TOP = 2.02
CONTENT_BOTTOM = 6.90
FIG_BOTTOM = 5.92


def _tf(box):
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def _emph_runs(p, text, size, color, bold=False):
    """Split on **...** and render those spans bold in the accent colour."""
    for i, part in enumerate(re.split(r"\*\*(.+?)\*\*", text)):
        if not part:
            continue
        r = p.add_run()
        r.text = part
        r.font.size = Pt(size)
        r.font.name = BODY_FONT
        strong = (i % 2 == 1)
        r.font.bold = bold or strong
        r.font.color.rgb = ACCENT if strong else color


def add_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(1, 0, 0, Emu(int(SW * 914400)), Emu(int(SH * 914400)))
    # Diagonal wash rather than flat white: adds depth without touching contrast.
    bg.fill.gradient()
    bg.fill.gradient_angle = 45.0
    stops = bg.fill.gradient_stops
    stops[0].color.rgb = BG_TOP
    stops[0].position = 0.0
    stops[1].color.rgb = BG_BOT
    stops[1].position = 1.0
    bg.line.fill.background()
    bg.shadow.inherit = False
    return s


def chrome(s, kicker, title, number):
    if kicker:
        b = s.shapes.add_textbox(Inches(ML), Inches(0.36), Inches(CW), Inches(0.32))
        p = _tf(b).paragraphs[0]
        r = p.add_run()
        r.text = kicker.upper()
        r.font.size = Pt(T_KICKER)
        r.font.bold = True
        r.font.name = BODY_FONT
        r.font.color.rgb = ACCENT
    ln = s.shapes.add_shape(1, Inches(ML), Inches(0.755), Inches(1.7), Inches(0.045))
    ln.fill.solid()
    ln.fill.fore_color.rgb = ACCENT
    ln.line.fill.background()
    ln.shadow.inherit = False
    b = s.shapes.add_textbox(Inches(ML), Inches(0.93), Inches(CW), Inches(1.02))
    tf = _tf(b)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.line_spacing = 1.06
    r = p.add_run()
    r.text = title
    r.font.size = Pt(T_TITLE)
    r.font.bold = True
    r.font.name = BODY_FONT
    r.font.color.rgb = DEEP
    if number:
        b = s.shapes.add_textbox(Inches(SW - 1.35), Inches(SH - 0.62), Inches(0.9), Inches(0.32))
        p = _tf(b).paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = str(number)
        r.font.size = Pt(T_PAGE)
        r.font.name = BODY_FONT
        r.font.color.rgb = PAGE


def bullets(s, items, top, width, left=ML, size=T_BULLET, gap=0.30):
    b = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width),
                             Inches(CONTENT_BOTTOM - top))
    tf = _tf(b)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.18
        p.space_after = Pt(gap * 72 * 0.45)
        dot = p.add_run()
        dot.text = "— "
        dot.font.size = Pt(size)
        dot.font.name = BODY_FONT
        dot.font.color.rgb = ACCENT
        _emph_runs(p, item, size, INK)
    return b


def picture(s, path, top, max_h, max_w=CW, left=None):
    w, h = Image.open(ROOT / path).size
    ar = w / h
    pw, ph = max_w, max_w / ar
    if ph > max_h:
        ph, pw = max_h, max_h * ar
    lx = left if left is not None else (SW - pw) / 2
    pic = s.shapes.add_picture(str(ROOT / path), Inches(lx), Inches(top),
                               Inches(pw), Inches(ph))
    # The figures are white-background PNGs; on the gradient wash a hairline border
    # makes them read as deliberate cards rather than floating rectangles.
    pic.line.color.rgb = RULE
    pic.line.width = Pt(0.75)
    return top + ph


def table(s, spec, top, dense=False):
    heads, rows, widths = spec["headers"], spec["rows"], spec["widths"]
    fs = 16 if dense else T_TABLE
    rh = min(0.44, (CONTENT_BOTTOM - top) / (len(rows) + 1)) if dense else 0.46
    total = sum(widths)
    left = (SW - total) / 2
    gf = s.shapes.add_table(len(rows) + 1, len(heads), Inches(left), Inches(top),
                            Inches(total), Inches(rh * (len(rows) + 1)))
    tbl = gf.table
    tbl.first_row = True
    tbl.horz_banding = False
    for j, w in enumerate(widths):
        tbl.columns[j].width = Inches(w)
    for i in range(len(rows) + 1):
        tbl.rows[i].height = Inches(rh)
    for j, htxt in enumerate(heads):
        c = tbl.cell(0, j)
        c.text = ""
        c.fill.solid()
        c.fill.fore_color.rgb = DEEP
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.margin_left = c.margin_right = Inches(0.10)
        p = c.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = htxt
        r.font.size = Pt(fs)
        r.font.bold = True
        r.font.name = BODY_FONT
        r.font.color.rgb = WHITE
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.text = ""
            c.fill.solid()
            c.fill.fore_color.rgb = PANEL if i % 2 else WHITE
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_left = c.margin_right = Inches(0.10)
            p = c.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT
            mark = "←" in val
            r = p.add_run()
            r.text = val
            r.font.size = Pt(fs)
            r.font.name = BODY_FONT
            r.font.bold = mark
            r.font.color.rgb = ACCENT if mark else INK
    return top + rh * (len(rows) + 1)


def code(s, lines, top):
    h = 0.40 * len(lines) + 0.42
    box = s.shapes.add_shape(1, Inches(ML), Inches(top), Inches(CW), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = PANEL
    box.line.color.rgb = RULE
    box.shadow.inherit = False
    tf = _tf(box)
    tf.margin_left = Inches(0.26)
    tf.margin_top = Inches(0.21)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.34
        for k, part in enumerate(re.split(r"\*\*(.+?)\*\*", line)):
            if not part:
                continue
            r = p.add_run()
            r.text = part
            r.font.size = Pt(T_CODE)
            r.font.name = MONO_FONT
            strong = (k % 2 == 1)
            r.font.bold = strong
            r.font.color.rgb = ACCENT if strong else INK
    return top + h


PHASES = [("Setup", "§1–3"), ("Data", "§4–7"), ("Exploration", "§8"), ("Features", "§9–10"),
          ("Models", "§11–13"), ("Decision", "§14–15"), ("Assurance", "§16–18")]


def diagram(s):
    n = len(PHASES)
    bw, gap = 1.50, 0.22
    total = n * bw + (n - 1) * gap
    x0 = (SW - total) / 2
    top, bh = 2.50, 1.25
    for i, (name, sec) in enumerate(PHASES):
        x = x0 + i * (bw + gap)
        box = s.shapes.add_shape(5, Inches(x), Inches(top), Inches(bw), Inches(bh))
        box.fill.solid()
        box.fill.fore_color.rgb = PANEL
        box.line.color.rgb = ACCENT
        box.line.width = Pt(1.5)
        box.shadow.inherit = False
        tf = _tf(box)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = Inches(0.05)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = 1.1
        r = p.add_run()
        r.text = name
        r.font.size = Pt(T_BULLET)
        r.font.bold = True
        r.font.name = BODY_FONT
        r.font.color.rgb = DEEP
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = sec
        r2.font.size = Pt(T_KICKER)
        r2.font.name = BODY_FONT
        r2.font.color.rgb = MUTED
        if i < n - 1:
            ar = s.shapes.add_shape(13, Inches(x + bw + 0.045),
                                    Inches(top + bh / 2 - 0.075),
                                    Inches(0.13), Inches(0.15))
            ar.fill.solid()
            ar.fill.fore_color.rgb = ACCENT
            ar.line.fill.background()
            ar.shadow.inherit = False
    band = s.shapes.add_shape(1, Inches(x0), Inches(4.28), Inches(total), Inches(0.78))
    band.fill.solid()
    band.fill.fore_color.rgb = DEEP
    band.line.fill.background()
    band.shadow.inherit = False
    tf = _tf(band)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "One rule throughout: use only what a buyer knew on the transaction date"
    r.font.size = Pt(T_BULLET)
    r.font.bold = True
    r.font.name = BODY_FONT
    r.font.color.rgb = WHITE
    b = s.shapes.add_textbox(Inches(x0), Inches(5.30), Inches(total), Inches(1.0))
    tf = _tf(b)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.2
    r = p.add_run()
    r.text = ("Four raw sources  →  59,946 × 37 master table  →  28 features  →  "
              "60/15/10/15 chronological split  →  14 models  →  conformal floor")
    r.font.size = Pt(T_SUB)
    r.font.name = BODY_FONT
    r.font.color.rgb = MUTED


def joingraph(s, rows):
    """Spine at the top, one labelled row per join, master table at the foot."""
    x0, w = 0.95, 11.45
    def band(top, h, txt, sub, fill, fg, bold=True):
        b = s.shapes.add_shape(5, Inches(x0), Inches(top), Inches(w), Inches(h))
        b.fill.solid(); b.fill.fore_color.rgb = fill
        b.line.color.rgb = ACCENT; b.line.width = Pt(1.25); b.shadow.inherit = False
        tf = _tf(b); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = Inches(0.16)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = txt
        r.font.size = Pt(T_BULLET); r.font.bold = bold
        r.font.name = BODY_FONT; r.font.color.rgb = fg
        if sub:
            r2 = p.add_run(); r2.text = "   " + sub
            r2.font.size = Pt(T_KICKER); r2.font.name = MONO_FONT; r2.font.color.rgb = fg
    band(2.02, 0.52, "Price history", "418,201 sales  ·  the spine", PANEL, DEEP)
    top = 2.70
    for src, mech, lag in rows:
        a = s.shapes.add_shape(13, Inches(x0 + 0.42), Inches(top + 0.055),
                               Inches(0.15), Inches(0.17))
        a.rotation = 90
        a.fill.solid(); a.fill.fore_color.rgb = ACCENT
        a.line.fill.background(); a.shadow.inherit = False
        b = s.shapes.add_textbox(Inches(x0 + 0.78), Inches(top), Inches(w - 0.78), Inches(0.30))
        tf = _tf(b); p = tf.paragraphs[0]
        r = p.add_run(); r.text = src + "  "
        r.font.size = Pt(T_SUB); r.font.bold = True
        r.font.name = BODY_FONT; r.font.color.rgb = INK
        r2 = p.add_run(); r2.text = mech
        r2.font.size = Pt(T_KICKER); r2.font.name = MONO_FONT; r2.font.color.rgb = MUTED
        if lag:
            r3 = p.add_run(); r3.text = "   LAGGED"
            r3.font.size = Pt(T_KICKER); r3.font.bold = True
            r3.font.name = MONO_FONT; r3.font.color.rgb = ACCENT
        top += 0.42
    band(top + 0.10, 0.52, "Master table", "59,946 rows  ×  37 columns", DEEP, WHITE)


def title_slide(s, spec):
    bar = s.shapes.add_shape(1, 0, Inches(2.42), Inches(SW), Inches(0.075))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    bar.shadow.inherit = False
    b = s.shapes.add_textbox(Inches(ML), Inches(1.42), Inches(CW), Inches(0.95))
    p = _tf(b).paragraphs[0]
    r = p.add_run()
    r.text = spec["title"]
    r.font.size = Pt(50)
    r.font.bold = True
    r.font.name = BODY_FONT
    r.font.color.rgb = DEEP
    b = s.shapes.add_textbox(Inches(ML), Inches(2.85), Inches(CW - 0.9), Inches(1.5))
    tf = _tf(b)
    p = tf.paragraphs[0]
    p.line_spacing = 1.24
    r = p.add_run()
    r.text = spec["subtitle"]
    r.font.size = Pt(22)
    r.font.name = BODY_FONT
    r.font.color.rgb = INK
    for txt, top, size, col in ((spec.get("authors"), 4.30, T_BULLET, INK),
                                (spec.get("course"), 4.80, T_SUB, MUTED)):
        if not txt:
            continue
        b = s.shapes.add_textbox(Inches(ML), Inches(top), Inches(CW), Inches(0.42))
        p = _tf(b).paragraphs[0]
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(size)
        r.font.name = BODY_FONT
        r.font.color.rgb = col
    b = s.shapes.add_textbox(Inches(ML), Inches(5.42), Inches(CW), Inches(0.9))
    tf = _tf(b)
    p = tf.paragraphs[0]
    p.line_spacing = 1.35
    r = p.add_run()
    r.text = spec["meta"]
    r.font.size = Pt(T_SUB)
    r.font.name = BODY_FONT
    r.font.color.rgb = MUTED


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


def build():
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)
    specs = MAIN + APPENDIX
    for n, spec in enumerate(specs, start=1):
        s = add_slide(prs)
        lay = spec["layout"]
        if lay == "title":
            title_slide(s, spec)
        else:
            chrome(s, spec.get("kicker"), spec["title"], n)
        if lay == "diagram":
            diagram(s)
        elif lay == "joins":
            joingraph(s, spec["joins"])
            bullets(s, spec["bullets"], 6.30, CW, size=T_SUB, gap=0.075)
        elif lay == "bullets":
            bullets(s, spec["bullets"], CONTENT_TOP + 0.22, CW)
        elif lay == "figure_below":
            nb = len(spec["bullets"])
            btop = CONTENT_BOTTOM - (0.305 * nb + 0.06)
            bot = picture(s, spec["figure"], CONTENT_TOP, btop - CONTENT_TOP - 0.16)
            bullets(s, spec["bullets"], max(btop, bot + 0.22), CW, size=T_SUB, gap=0.075)
        elif lay == "figure_full":
            picture(s, spec["figure"], CONTENT_TOP + 0.15, CONTENT_BOTTOM - CONTENT_TOP - 0.15)
        elif lay == "table":
            end = table(s, spec["table"], CONTENT_TOP + 0.12, dense=spec.get("dense", False))
            if spec.get("bullets"):
                bullets(s, spec["bullets"], end + 0.30, CW, size=T_SUB, gap=0.18)
        elif lay == "code":
            end = code(s, spec["code"], CONTENT_TOP + 0.10)
            bullets(s, spec["bullets"], end + 0.36, CW, size=T_SUB, gap=0.22)
        notes(s, spec["notes"])
    out = ROOT / "presentation.pptx"
    prs.save(str(out))
    print(f"wrote {out}  ({len(specs)} slides: {len(MAIN)} main + {len(APPENDIX)} appendix)")


if __name__ == "__main__":
    build()
